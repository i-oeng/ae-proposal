from __future__ import annotations

import json
import logging
import os
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Any

from pydantic import ValidationError
from pptx import Presentation

from core.anthropic_client import get_anthropic_client
from core.config_loader import AppConfig
from core.extraction import DEFAULT_ANTHROPIC_MODEL, _content_block_for_file, _extract_json_object
from core.extraction_cache import cache_key_for_files, load_cached_model, store_cached_model
from core.models import ClientInfoDraft
from core.utils import load_local_env

logger = logging.getLogger(__name__)
CLIENT_EXTRACTION_CACHE_VERSION = "client-extraction-v1"

SUPPORTED_INDUSTRIES = {
    "manufacturing",
    "cold_storage",
    "food_processing",
    "retail",
    "hospitality",
}

COUNTRY_ALIASES = {
    "ghana": "Ghana",
    "nigeria": "Nigeria",
    "senegal": "Senegal",
    "cote d'ivoire": "Cote d'Ivoire",
    "côte d'ivoire": "Cote d'Ivoire",
    "cote divoire": "Cote d'Ivoire",
    "ivory coast": "Cote d'Ivoire",
}


def _normalize_country(value: str | None) -> str | None:
    if not value:
        return None
    return COUNTRY_ALIASES.get(value.strip().lower(), value.strip())


def _normalize_industry(value: str | None) -> str | None:
    if not value:
        return None
    normalized = value.strip().lower().replace("-", "_").replace(" ", "_")
    if normalized in SUPPORTED_INDUSTRIES:
        return normalized
    if "food" in normalized:
        return "food_processing"
    if "cold" in normalized or "storage" in normalized or "refriger" in normalized:
        return "cold_storage"
    if "manufact" in normalized or "factory" in normalized:
        return "manufacturing"
    if "retail" in normalized or "shop" in normalized or "store" in normalized:
        return "retail"
    if "hotel" in normalized or "hospitality" in normalized:
        return "hospitality"
    return normalized


def _extract_pdf_text(file_path: str) -> str:
    try:
        import pdfplumber
    except ImportError:
        return ""

    lines: list[str] = []
    try:
        with pdfplumber.open(file_path) as pdf:
            for page_index, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    lines.append(f"Page {page_index}:\n{text.strip()}")
    except Exception as exc:  # noqa: BLE001
        logger.warning("PDF text extraction failed for %s: %s", file_path, exc)
        return ""
    return "\n\n".join(lines)


def _extract_pptx_text(file_path: str) -> str:
    prs = Presentation(file_path)
    lines: list[str] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    slide_lines.append(text)
        if slide_lines:
            lines.append(f"Slide {slide_index}:\n" + "\n".join(slide_lines))
    return "\n\n".join(lines)


def _content_block_for_client_file(file_path: str) -> dict[str, Any]:
    path = Path(file_path)
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        text = _extract_pptx_text(file_path)
        return {
            "type": "text",
            "text": (
                f"Extracted text from uploaded PowerPoint file {path.name}:\n\n"
                f"{text or '[No text found in PowerPoint shapes.]'}"
            ),
        }
    if suffix == ".pdf":
        text = _extract_pdf_text(file_path)
        if len(text.strip()) >= 80:
            return {
                "type": "text",
                "text": (
                    f"Extracted text from uploaded PDF file {path.name}:\n\n"
                    f"{text[:45000]}"
                ),
            }
    return _content_block_for_file(file_path)


def _content_blocks_for_client_files(file_paths: list[str]) -> list[dict[str, Any]]:
    if len(file_paths) < 2:
        return [_content_block_for_client_file(path) for path in file_paths]
    worker_count = min(len(file_paths), 4)
    with ThreadPoolExecutor(max_workers=worker_count) as executor:
        return list(executor.map(_content_block_for_client_file, file_paths))


def _call_claude_for_client_info(file_paths: list[str]) -> dict[str, Any]:
    load_local_env()
    client = get_anthropic_client()
    model = os.getenv("ANTHROPIC_VISION_MODEL", DEFAULT_ANTHROPIC_MODEL)
    prompt = """
Extract client information for a commercial solar/PPA proposal.

Use the uploaded client documents, site visit reports, technical feasibility reports, prior commercial offers,
presentations, PDFs, screenshots, or images.
Return strict JSON only. Use null when a value is not present and cannot be reasonably inferred.
Reasonable assumptions are allowed, but every assumption must be listed in extraction_notes.

Extraction guidance:
- Business descriptions may explicitly state the company activity; map cocoa, beverage, food factory,
  cold rooms, roasting, processing, packaging, and agro-industry to food_processing unless a better
  supported industry is clearly stated.
- Site visit reports may contain fields such as "Site", "Client", "Coordonnees GPS", "transformateur",
  "groupe electrogene", "Surface Totale", "Puissance Estimee", and roof/building tables. Use those
  to infer city, coordinates, grid_connection_kva, has_diesel_generators, and available_roof_area_m2.
- Convert GPS coordinates from degrees/minutes/seconds to decimal latitude and longitude.
- If a report lists multiple roof areas, use the sum of usable or listed roof areas for available_roof_area_m2
  and explain the basis in extraction_notes. Do not use proposed kWp as roof area.
- Existing offer decks can confirm client name, country, project size, PPA term, tariff references, and
  diesel assumptions, but do not overwrite stronger site-report values with marketing summary values.
- If a proposed ASPAN/PPA solar tariff per kWh is explicitly stated, put it in ppa_tariff_per_kwh_override.
- If a diesel fuel price per liter is explicitly stated, put it in diesel_price_per_liter_override.
- Prefer source values from technical reports over older commercial offer slides when there is a conflict.

Supported industry values:
manufacturing, cold_storage, food_processing, retail, hospitality

Supported countries:
Ghana, Nigeria, Senegal, Cote d'Ivoire

Return exactly this JSON object:
{
  "client_name": string | null,
  "industry": one supported industry string | null,
  "country": one supported country string | null,
  "city": string | null,
  "latitude": number | null,
  "longitude": number | null,
  "business_description": string | null,
  "has_diesel_generators": true | false | null,
  "grid_connection_kva": number | null,
  "available_roof_area_m2": number | null,
  "daytime_fraction_override": number between 0 and 1 | null,
  "ppa_tariff_per_kwh_override": number | null,
  "diesel_price_per_liter_override": number | null,
  "extraction_notes": ["short note"],
  "field_confidence": {
    "client_name": number between 0 and 1,
    "industry": number between 0 and 1,
    "country": number between 0 and 1,
    "business_description": number between 0 and 1,
    "has_diesel_generators": number between 0 and 1,
    "grid_connection_kva": number between 0 and 1,
    "available_roof_area_m2": number between 0 and 1
  }
}
"""
    content = [{"type": "text", "text": prompt}, *_content_blocks_for_client_files(file_paths)]
    message = client.messages.create(
        model=model,
        max_tokens=1600,
        messages=[{"role": "user", "content": content}],
    )
    text = "".join(block.text for block in message.content if getattr(block, "type", None) == "text")
    return _extract_json_object(text)


def fallback_client_info(reason: str) -> ClientInfoDraft:
    return ClientInfoDraft(
        extraction_notes=[reason, "Client information was not extracted; review fields manually."],
        field_confidence={},
    )


def extract_client_info(file_paths: list[str], config: AppConfig) -> ClientInfoDraft:  # noqa: ARG001
    load_local_env()
    existing_paths = [path for path in file_paths if Path(path).exists()]
    if not existing_paths:
        return fallback_client_info("No client information files uploaded.")
    if not os.getenv("ANTHROPIC_API_KEY"):
        return fallback_client_info("ANTHROPIC_API_KEY not set; client extraction fallback used.")

    model = os.getenv("ANTHROPIC_VISION_MODEL", DEFAULT_ANTHROPIC_MODEL)
    cache_key = cache_key_for_files(CLIENT_EXTRACTION_CACHE_VERSION, existing_paths, model)
    cached = load_cached_model(cache_key, ClientInfoDraft)
    if cached is not None:
        return cached

    last_error: Exception | None = None
    for _attempt in range(2):
        try:
            payload = _call_claude_for_client_info(existing_paths)
            payload["country"] = _normalize_country(payload.get("country"))
            payload["industry"] = _normalize_industry(payload.get("industry"))
            draft = ClientInfoDraft.model_validate(payload)
            store_cached_model(cache_key, draft)
            return draft
        except (json.JSONDecodeError, ValidationError, OSError, Exception) as exc:  # noqa: BLE001
            logger.warning("Client info extraction attempt failed: %s", exc)
            last_error = exc

    return fallback_client_info(
        f"Claude client extraction failed twice using model "
        f"{os.getenv('ANTHROPIC_VISION_MODEL', DEFAULT_ANTHROPIC_MODEL)}: {last_error}"
    )
