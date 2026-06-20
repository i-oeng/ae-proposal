from __future__ import annotations

from datetime import date
from pathlib import Path
import re

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from core.config_loader import AppConfig
from core.models import BillData, CalcResult, ClientInfo, NarrativeSections
from core.utils import ensure_dir, format_currency, format_kwp, format_number, safe_slug


TITLE_COLOR = "173B33"
ACCENT_COLOR = "00A884"
AMBER_COLOR = "F4B942"
DARK_COLOR = "1F2933"
MUTED_COLOR = "6B7280"
LIGHT_COLOR = "F7FAF8"


def ensure_placeholder_template(path: str | Path = "templates/aspan_template.pptx") -> Path:
    template_path = Path(path)
    template_path.parent.mkdir(parents=True, exist_ok=True)
    if template_path.exists():
        return template_path

    prs = Presentation()
    prs.core_properties.title = "Aspan Energy placeholder proposal template"
    prs.core_properties.subject = (
        "Editable Aspan Energy proposal template. Brand assets and final commercial "
        "terms should be confirmed before client use."
    )
    prs.save(template_path)
    return template_path


def _rgb(hex_color: str):
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(hex_color)


def _add_title(slide, title: str, subtitle: str | None = None) -> None:
    shape = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(9.1), Inches(0.55))
    frame = shape.text_frame
    frame.clear()
    run = frame.paragraphs[0].add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = _rgb(TITLE_COLOR)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.88), Inches(8.8), Inches(0.32))
        sub_frame = sub.text_frame
        sub_frame.text = subtitle
        sub_frame.paragraphs[0].runs[0].font.size = Pt(9)
        sub_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(MUTED_COLOR)


def _add_band(slide) -> None:
    band = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.12))
    band.fill.solid()
    band.fill.fore_color.rgb = _rgb(ACCENT_COLOR)
    band.line.fill.background()


def _add_text_box(slide, x: float, y: float, w: float, h: float, text: str, size: int = 13) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.text = text
    for paragraph in frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = _rgb(DARK_COLOR)


def _add_section_heading(slide, x: float, y: float, w: float, text: str) -> None:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.28))
    frame = shape.text_frame
    frame.text = text
    paragraph = frame.paragraphs[0]
    paragraph.runs[0].font.size = Pt(10)
    paragraph.runs[0].font.bold = True
    paragraph.runs[0].font.color.rgb = _rgb(TITLE_COLOR)


def _add_bullets(slide, x: float, y: float, w: float, h: float, bullets: list[str], size: int = 12) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    for index, item in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = _rgb(DARK_COLOR)


def _add_metric(slide, x: float, y: float, label: str, value: str, width: float = 2.1) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(width), Inches(0.78))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(LIGHT_COLOR)
    shape.line.color.rgb = _rgb(ACCENT_COLOR)
    frame = shape.text_frame
    frame.clear()
    label_p = frame.paragraphs[0]
    label_p.text = label.upper()
    label_p.font.size = Pt(7)
    label_p.font.color.rgb = _rgb(MUTED_COLOR)
    value_p = frame.add_paragraph()
    value_p.text = value
    value_p.font.size = Pt(15)
    value_p.font.bold = True
    value_p.font.color.rgb = _rgb(TITLE_COLOR)


def _new_slide(prs: Presentation, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_band(slide)
    _add_title(slide, title, subtitle)
    return slide


def _clean_text(text: str | None) -> str:
    if not text:
        return ""
    cleaned = re.sub(r"\b(?:TODO|TBD)\b:?", "", text, flags=re.IGNORECASE)
    return re.sub(r"\s+", " ", cleaned).strip()


def _is_placeholder_text(text: str) -> bool:
    lowered = text.lower()
    return not text or lowered in {
        "details to be confirmed by aspan energy.",
        "details to be confirmed by aspan energy",
    }


def _text_to_bullets(text: str | None, fallback: list[str], limit: int = 4) -> list[str]:
    cleaned = _clean_text(text)
    if _is_placeholder_text(cleaned):
        return fallback[:limit]
    parts = re.split(r"(?<=[.!?])\s+|;\s+|\n+", cleaned)
    bullets = [part.strip(" .") for part in parts if len(part.strip()) > 12]
    return (bullets or fallback)[:limit]


def _current_energy_context(bill: BillData, client: ClientInfo) -> str:
    diesel_line = (
        "The presence of diesel generation makes avoided fuel use a separate value driver."
        if client.has_diesel_generators
        else "Diesel generation has not been confirmed, so the grid-only case should remain the base case."
    )
    return (
        "The reviewed bills indicate a material industrial electricity load with exposure to both "
        f"active energy tariff and non-energy charges. {diesel_line} Final commercial use should "
        "confirm billing periods, penalties, contracted demand, and any unpaid balances excluded from "
        "the monthly cost baseline."
    )


def _solution_context(client: ClientInfo, calc: CalcResult) -> str:
    return (
        "The proposed PV size is constrained by the most conservative sizing input rather than by a "
        "single headline project target. For this run, the binding constraint is "
        f"{calc.pv_recommendation.binding_constraint.replace('_', ' ')}. Final sizing should be "
        "reconfirmed after roof layout, structural access, interconnection, and operating schedule review."
    )


def _client_context_bullets(client: ClientInfo, calc: CalcResult) -> list[str]:
    bullets = [
        f"Industry: {client.industry.replace('_', ' ').title()}",
        f"Country: {client.country}",
        f"Diesel generators: {'confirmed' if client.has_diesel_generators else 'not confirmed'}",
    ]
    if client.business_description:
        bullets.insert(0, f"Client activity: {client.business_description}")
    if client.city:
        bullets.append(f"Site location: {client.city}")
    if client.grid_connection_kva is not None:
        bullets.append(f"Grid connection capacity: {format_number(client.grid_connection_kva)} kVA")
    if client.available_roof_area_m2 is not None:
        bullets.append(f"Available roof area: {format_number(client.available_roof_area_m2)} m2")
    bullets.append(
        f"Sizing basis: {calc.pv_recommendation.binding_constraint.replace('_', ' ')} constraint"
    )
    return bullets[:6]


def _market_context_bullets(client: ClientInfo) -> list[str]:
    bullets: list[str] = []
    if client.country.strip().lower() == "cote d'ivoire":
        bullets.append(
            "Cote d'Ivoire context: ICCO is located in Abidjan, reinforcing the country's role in the cocoa economy."
        )
    if client.industry.strip().lower() == "food_processing":
        bullets.append(
            "Food and cocoa processing can include roasting, grinding, pressing, cooling, packaging, and process utility loads."
        )
    bullets.append(
        "World Bank Global Solar Atlas is useful for preliminary solar-resource screening; final production remains site-specific."
    )
    bullets.append(
        "Client-specific values in this deck come from uploaded bills, site reports, and offer documents rather than public web data."
    )
    return bullets


def _market_source_note(client: ClientInfo) -> str:
    sources = ["Client documents"]
    if client.country.strip().lower() == "cote d'ivoire" or client.industry.strip().lower() == "food_processing":
        sources.append("ICCO")
    sources.append("World Bank Global Solar Atlas")
    return "Source notes: " + "; ".join(sources) + "."


def _delivery_scope_bullets(text: str | None) -> list[str]:
    return _text_to_bullets(
        text,
        [
            "Review customer energy data, bills, tariff structure, and diesel assumptions",
            "Coordinate site assessment inputs and preliminary technical feasibility review",
            "Refine PV sizing, savings analysis, and commercial proposal assumptions",
            "Confirm final engineering, legal, interconnection, O&M, and performance obligations",
        ],
        limit=4,
    )


def _delivery_process_bullets(text: str | None) -> list[str]:
    return _text_to_bullets(
        text,
        [
            "Confirm utility bill, tariff, and diesel baseline assumptions",
            "Validate roof area, usable layout, structural access, and grid connection information",
            "Complete technical feasibility and commercial review",
            "Prepare detailed proposal and implementation plan for approval",
        ],
        limit=4,
    )


def _next_step_bullets(text: str | None) -> list[str]:
    return _text_to_bullets(
        text,
        [
            "Confirm assumptions and commercial review owners",
            "Validate roof area and site constraints",
            "Confirm tariff, diesel, and PPA tariff inputs",
            "Proceed to detailed proposal after review",
        ],
        limit=4,
    )


def _configure_savings_chart(chart, currency: str) -> None:
    chart.has_title = True
    chart.chart_title.text_frame.text = f"Savings comparison ({currency})"
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = f"Savings ({currency})"
    chart.value_axis.tick_labels.number_format = f'"{currency}" #,##0'
    chart.category_axis.has_title = True
    chart.category_axis.axis_title.text_frame.text = "Savings period"
    chart.has_legend = True


def _scenario_bullets(result, currency: str) -> list[str]:
    return [
        f"Monthly savings year 1: {format_currency(result.monthly_savings_year_1, currency)}",
        f"Annual savings year 1: {format_currency(result.annual_savings_year_1, currency)}",
        f"Cumulative savings: {format_currency(result.cumulative_savings, currency)}",
        f"Year 1 solar used: {format_number(result.year_1_solar_used_kwh)} kWh",
        f"Current cost baseline: {format_currency(result.current_cost_per_kwh_year_1, currency, 3)} per kWh",
        f"PPA tariff: {format_currency(result.ppa_tariff_per_kwh, currency, 3)} per kWh",
    ]


def build_pptx(
    bill: BillData,
    client: ClientInfo,
    calc: CalcResult,
    narrative: NarrativeSections,
    config: AppConfig,
    template_path: str | Path = "templates/aspan_template.pptx",
) -> str:
    ensure_placeholder_template(template_path)
    prs = Presentation(str(template_path))
    output_dir = ensure_dir(config.output_dir)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_band(slide)
    _add_text_box(slide, 0.65, 0.75, 8.4, 0.7, "Aspan Energy", 28)
    _add_text_box(slide, 0.68, 1.38, 8.8, 0.4, "Commercial solar PPA analysis", 14)
    _add_text_box(
        slide,
        0.68,
        2.25,
        7.8,
        1.7,
        f"{client.client_name}\n{client.country} | {client.industry.replace('_', ' ').title()}\nProposal date: {date.today().isoformat()}",
        20,
    )
    _add_metric(slide, 0.72, 5.15, "PV size", format_kwp(calc.pv_recommendation.recommended_kwp), 2.0)
    _add_metric(slide, 2.95, 5.15, "PPA tariff", format_currency(calc.ppa_tariff_per_kwh, bill.currency, 3), 2.0)
    _add_metric(slide, 5.18, 5.15, "Analysis", f"{config.analysis_horizon_years} years", 2.0)
    _add_metric(slide, 7.41, 5.15, "Currency", bill.currency, 1.6)
    _add_text_box(
        slide,
        0.72,
        6.15,
        8.6,
        0.45,
        "Preliminary proposal generated from reviewed documents; final engineering and commercial terms require Aspan confirmation.",
        9,
    )

    # Slide 2
    slide = _new_slide(prs, "Executive Summary")
    _add_metric(slide, 0.65, 1.35, "PV size", format_kwp(calc.pv_recommendation.recommended_kwp))
    _add_metric(slide, 2.95, 1.35, "Annual production", f"{format_number(calc.annual_solar_production_kwh)} kWh")
    _add_metric(slide, 5.25, 1.35, "Grid annual savings", format_currency(calc.scenario_grid_replacement.annual_savings_year_1, bill.currency))
    _add_metric(slide, 7.55, 1.35, "Diesel annual savings", format_currency(calc.scenario_grid_diesel.annual_savings_year_1, bill.currency))
    _add_text_box(slide, 0.72, 2.45, 8.6, 1.45, narrative.executive_summary, 13)
    _add_bullets(
        slide,
        0.72,
        4.15,
        8.8,
        1.1,
        [
            f"Grid Replacement cumulative savings: {format_currency(calc.scenario_grid_replacement.cumulative_savings, bill.currency)}",
            f"Grid + Diesel cumulative savings: {format_currency(calc.scenario_grid_diesel.cumulative_savings, bill.currency)}",
        ],
    )
    _add_bullets(slide, 0.72, 5.55, 8.8, 0.85, _client_context_bullets(client, calc)[:3], 9)

    # Slide 3
    slide = _new_slide(prs, "Current Energy Situation")
    _add_bullets(
        slide,
        0.7,
        1.35,
        4.1,
        2.3,
        [
            f"Monthly consumption: {format_number(bill.monthly_kwh)} kWh",
            f"Monthly electricity cost: {format_currency(bill.total_cost, bill.currency)}",
            f"Grid tariff: {format_currency(bill.tariff_per_kwh, bill.currency, 3)} per kWh",
            f"Tariff basis: {bill.tariff_basis.replace('_', ' ')}",
            f"Diesel generators: {'Yes' if client.has_diesel_generators else 'No'}",
            *(
                [f"Penalties: {format_currency(bill.penalties, bill.currency)}"]
                if bill.penalties is not None
                else []
            ),
        ],
    )
    current_text = _clean_text(narrative.current_energy_situation)
    if _is_placeholder_text(current_text):
        current_text = _current_energy_context(bill, client)
    _add_text_box(slide, 5.05, 1.35, 4.2, 2.8, current_text, 12)

    # Slide 4
    slide = _new_slide(prs, "Proposed Solar Solution")
    _add_bullets(
        slide,
        0.7,
        1.35,
        4.3,
        2.6,
        [
            f"Recommended size: {format_kwp(calc.pv_recommendation.recommended_kwp)}",
            f"Binding constraint: {calc.pv_recommendation.binding_constraint}",
            f"Estimated annual production: {format_number(calc.annual_solar_production_kwh)} kWh",
            f"Daytime consumption offset cap: {format_number(calc.annual_daytime_consumption_kwh)} kWh",
        ],
    )
    solution_text = _clean_text(narrative.proposed_solar_solution)
    if _is_placeholder_text(solution_text) or len(solution_text) < 120:
        solution_text = _solution_context(client, calc)
    _add_text_box(slide, 5.1, 1.35, 4.1, 2.8, solution_text, 12)

    # Slide 5
    slide = _new_slide(prs, "PPA Model Explanation")
    _add_text_box(slide, 0.72, 1.35, 8.7, 1.65, narrative.ppa_model_explanation, 14)
    _add_bullets(
        slide,
        0.72,
        3.35,
        8.7,
        1.45,
        [
            "No upfront CAPEX, subject to final commercial confirmation.",
            "Client pays for solar energy through a PPA.",
            "Solar reduces daytime grid or diesel reliance.",
            "Contract details must be confirmed by Aspan Energy.",
        ],
    )

    # Slide 6
    slide = _new_slide(prs, "Economic Analysis: Grid Replacement")
    _add_bullets(slide, 0.72, 1.35, 4.5, 2.1, _scenario_bullets(calc.scenario_grid_replacement, bill.currency))
    econ_text = _clean_text(narrative.economic_analysis_summary)
    if _is_placeholder_text(econ_text):
        econ_text = (
            "This scenario isolates savings where solar replaces grid electricity priced at the reviewed "
            "tariff baseline. The cumulative case applies the configured tariff escalation and PV degradation "
            "assumptions across the analysis horizon."
        )
    _add_text_box(slide, 5.45, 1.35, 3.8, 2.4, econ_text, 12)

    # Slide 7
    slide = _new_slide(prs, "Economic Analysis: Grid + Diesel")
    _add_bullets(slide, 0.72, 1.35, 4.5, 2.1, _scenario_bullets(calc.scenario_grid_diesel, bill.currency))
    _add_text_box(
        slide,
        5.45,
        1.35,
        3.8,
        2.4,
        "Scenario assumes a 50% grid and 50% diesel operating baseline for displaced daytime energy.",
        13,
    )

    # Slide 8
    slide = _new_slide(
        prs,
        "Scenario Comparison",
        f"Values shown in {bill.currency}; comparison uses year 1 annual savings and {config.analysis_horizon_years}-year cumulative savings.",
    )
    chart_data = ChartData()
    chart_data.categories = ["Annual savings", f"{config.analysis_horizon_years}-year savings"]
    chart_data.add_series(
        "Grid Replacement",
        (
            calc.scenario_grid_replacement.annual_savings_year_1,
            calc.scenario_grid_replacement.cumulative_savings,
        ),
    )
    chart_data.add_series(
        "Grid + Diesel",
        (
            calc.scenario_grid_diesel.annual_savings_year_1,
            calc.scenario_grid_diesel.cumulative_savings,
        ),
    )
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8),
        Inches(1.35),
        Inches(8.4),
        Inches(3.9),
        chart_data,
    )
    _configure_savings_chart(chart_shape.chart, bill.currency)
    _add_bullets(
        slide,
        0.82,
        5.48,
        8.35,
        0.85,
        [
            f"Unit: {bill.currency}; values are not kWh and should be read as monetary savings.",
            "Grid + Diesel uses a blended current-cost baseline for displaced daytime energy.",
        ],
        9,
    )

    # Slide 9
    slide = _new_slide(prs, "Market Context & Proposal Basis")
    _add_bullets(slide, 0.72, 1.35, 8.55, 2.25, _market_context_bullets(client), 12)
    _add_text_box(slide, 0.72, 4.05, 8.6, 0.55, _market_source_note(client), 9)
    _add_bullets(slide, 0.72, 4.72, 8.55, 1.1, _client_context_bullets(client, calc), 9)

    # Slide 10
    slide = _new_slide(prs, "Delivery Plan & Next Steps")
    _add_section_heading(slide, 0.7, 1.28, 2.6, "Aspan scope")
    _add_bullets(slide, 0.7, 1.68, 2.65, 3.6, _delivery_scope_bullets(narrative.scope_of_services), 10)
    _add_section_heading(slide, 3.65, 1.28, 2.65, "Review process")
    _add_bullets(slide, 3.65, 1.68, 2.65, 3.6, _delivery_process_bullets(narrative.implementation_process), 10)
    _add_section_heading(slide, 6.55, 1.28, 2.65, "Immediate next steps")
    _add_bullets(slide, 6.55, 1.68, 2.75, 3.6, _next_step_bullets(narrative.next_steps), 10)
    _add_text_box(
        slide,
        0.72,
        5.55,
        8.55,
        0.65,
        "Final dates, responsibilities, contract terms, and performance obligations remain subject to Aspan and client confirmation.",
        9,
    )

    # Slide 11
    slide = _new_slide(prs, "Assumptions Appendix")
    assumption_lines = [
        f"Monthly kWh: {format_number(calc.assumptions.monthly_kwh)}",
        f"Daytime fraction: {calc.assumptions.daytime_fraction:.2f}",
        f"Roof area: {format_number(calc.assumptions.roof_area_m2 or 0)} m2",
        f"Grid capacity: {format_number(calc.assumptions.grid_capacity_kva or 0)} kVA",
        f"Recommended kWp: {format_kwp(calc.assumptions.recommended_kwp)}",
        f"PPA tariff: {format_currency(calc.assumptions.ppa_tariff_per_kwh, bill.currency, 3)} per kWh",
        f"Grid tariff: {format_currency(calc.assumptions.grid_tariff_per_kwh, bill.currency, 3)} per kWh",
        f"Tariff basis: {bill.tariff_basis.replace('_', ' ')}",
        f"Active energy charge: {format_currency(bill.active_energy_charge or 0, bill.currency)}",
        f"Penalties: {format_currency(bill.penalties or 0, bill.currency)}",
        f"Diesel price: {format_currency(calc.assumptions.diesel_price_per_liter, bill.currency, 2)} per liter",
        f"Diesel kWh per liter: {calc.assumptions.diesel_kwh_per_liter:.2f}",
        f"Performance ratio: {calc.assumptions.performance_ratio:.2f}",
        f"Panel degradation rate: {calc.assumptions.panel_degradation_rate:.3f}",
        f"Tariff escalation rate: {calc.assumptions.tariff_escalation_rate:.2f}",
        f"Analysis horizon: {calc.assumptions.analysis_horizon_years} years",
        f"Production source: {calc.assumptions.production_source}",
    ]
    assumption_lines.extend(calc.assumptions.confidence_notes)
    assumption_lines.extend(calc.assumptions.extraction_notes)
    _add_bullets(slide, 0.6, 1.18, 8.9, 4.75, assumption_lines[:18], 9)

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.LEFT

    filename = f"{safe_slug(client.client_name)}_{date.today().isoformat()}_proposal.pptx"
    output_path = output_dir / filename
    prs.save(output_path)
    return str(output_path)
