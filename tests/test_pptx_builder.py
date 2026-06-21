from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from core.calc_engine import calculate_proposal
from core.config_loader import AppConfig
from core.models import BillData, ClientInfo, NarrativeSections
from core.pptx_builder import build_pptx


def _config(tmp_path: Path) -> AppConfig:
    return AppConfig(
        analysis_horizon_years=15,
        kwp_per_m2=0.15,
        grid_capacity_solar_fraction=0.70,
        daytime_fraction_default=0.70,
        default_specific_yield_kwh_per_kwp=1450,
        conservative_roof_area_m2_default=1200,
        daytime_fraction_by_industry={"food_processing": 0.75, "default": 0.70},
        performance_ratio=0.80,
        panel_degradation_rate=0.005,
        use_nasa_power=False,
        tariff_escalation_rate=0.03,
        ppa_discount_to_grid_tariff=0.15,
        ppa_tariff_per_kwh_default=None,
        diesel_price_per_liter=700,
        diesel_kwh_per_liter=3.5,
        system_cost_per_kwp=850,
        default_latitude_by_country={"Cote d'Ivoire": 5.36},
        default_longitude_by_country={"Cote d'Ivoire": -4.0083},
        nasa_power_cache_dir=str(tmp_path / "nasa"),
        output_dir=str(tmp_path),
        proposal_log_path=str(tmp_path / "proposal_runs.jsonl"),
    )


def _all_slide_text(prs: Presentation) -> str:
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                texts.extend(paragraph.text for paragraph in shape.text_frame.paragraphs)
    return "\n".join(texts)


def test_build_pptx_has_no_todos_and_labels_savings_chart(tmp_path: Path) -> None:
    config = _config(tmp_path)
    bill = BillData(
        monthly_kwh=233633,
        currency="XOF",
        total_cost=23546304,
        tariff_per_kwh=75.078,
        penalties=2860036,
        active_energy_charge=17531792,
        tariff_basis="kwh_weighted_average_across_uploaded_bills",
        field_confidence={"monthly_kwh": 0.97},
    )
    client = ClientInfo(
        client_name="NESKAO",
        industry="food_processing",
        country="Cote d'Ivoire",
        city="Abidjan",
        business_description=(
            "Leading food and beverage manufacturing company specializing in the production and "
            "distribution of cocoa-based products and instant beverages for the West African market, "
            "operating industrial production facilities in Abidjan's Zone Industrielle."
        ),
        has_diesel_generators=True,
        grid_connection_kva=1250,
        available_roof_area_m2=5149,
    )
    calc = calculate_proposal(bill, client, config)
    narrative = NarrativeSections(
        executive_summary="Aspan Energy proposes a preliminary solar PPA solution for NESKAO.",
        current_energy_situation="Details to be confirmed by Aspan Energy.",
        proposed_solar_solution="The recommended system size should be validated against roof and grid constraints.",
        ppa_model_explanation="The client pays for generated solar energy under a PPA structure.",
        economic_analysis_summary="Details to be confirmed by Aspan Energy.",
        scope_of_services="Review energy data; coordinate site assessment; prepare detailed proposal.",
        implementation_process="Confirm assumptions; validate site constraints; complete commercial review.",
        next_steps="Confirm inputs; validate roof area; proceed to detailed proposal.",
    )

    output = build_pptx(bill, client, calc, narrative, config, template_path=tmp_path / "template.pptx")
    prs = Presentation(output)
    all_text = _all_slide_text(prs)

    assert "TODO" not in all_text
    assert "Market Context & Proposal Basis" in all_text
    assert "Delivery Plan & Next Steps" in all_text
    assert "Scope of Aspan Services" not in all_text
    assert "Implementation Process" not in all_text
    assert "Unit: XOF" in all_text

    charts = [
        shape.chart
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_chart", False)
    ]
    assert charts
    assert charts[0].value_axis.axis_title.text_frame.text == "Savings (XOF)"

    long_text_frames = [
        shape.text_frame
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and len(shape.text) > 80
    ]
    assert long_text_frames
    assert all(frame.word_wrap for frame in long_text_frames)

    delivery_slide = prs.slides[9]
    delivery_frames = [
        shape.text_frame
        for shape in delivery_slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    ]
    assert all(frame.word_wrap for frame in delivery_frames if len(frame.text) > 30)
