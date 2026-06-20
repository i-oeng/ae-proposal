from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


MIN_WORDS_BY_SLIDE = {
    1: 8,
    2: 35,
    3: 28,
    4: 28,
    5: 24,
    6: 24,
    7: 24,
    8: 12,
    9: 24,
    10: 24,
    11: 24,
}


def _shape_text(shape) -> str:
    if not hasattr(shape, "text"):
        return ""
    return str(shape.text or "").strip()


def _slide_text(slide) -> str:
    return "\n".join(text for text in (_shape_text(shape) for shape in slide.shapes) if text)


def _word_count(text: str) -> int:
    return len(re.findall(r"[A-Za-z0-9]+", text))


def _chart_titles(slide) -> list[str]:
    titles: list[str] = []
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.CHART:
            continue
        chart = shape.chart
        if chart.has_title:
            titles.append(chart.chart_title.text_frame.text)
        if chart.value_axis.has_title:
            titles.append(chart.value_axis.axis_title.text_frame.text)
    return titles


def qa_pptx(path: Path) -> int:
    if not path.exists():
        print(f"Missing PPTX: {path}")
        return 2

    prs = Presentation(path)
    failures: list[str] = []
    all_text = "\n".join(_slide_text(slide) for slide in prs.slides)

    if re.search(r"\b(?:TODO|TBD)\b", all_text, flags=re.IGNORECASE):
        failures.append("Deck contains TODO/TBD text.")

    for index, slide in enumerate(prs.slides, start=1):
        text = _slide_text(slide)
        word_count = _word_count(text)
        minimum = MIN_WORDS_BY_SLIDE.get(index, 12)
        if word_count < minimum:
            failures.append(f"Slide {index} has low text density: {word_count} words, expected at least {minimum}.")

    chart_titles = [title for slide in prs.slides for title in _chart_titles(slide)]
    if not chart_titles:
        failures.append("Deck has no native PowerPoint charts.")
    elif not any("XOF" in title or "savings" in title.lower() for title in chart_titles):
        failures.append("Chart titles/axis labels do not clearly include savings units.")

    print(f"Checked {path}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Chart labels: {', '.join(chart_titles) if chart_titles else 'none'}")
    if failures:
        print("QA failed:")
        for failure in failures:
            print(f"- {failure}")
        return 1

    print("QA passed: no TODO/TBD text, slide density is acceptable, and charts have labels.")
    return 0


if __name__ == "__main__":
    target = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("outputs/neskao_2026-06-20_proposal.pptx")
    raise SystemExit(qa_pptx(target))
